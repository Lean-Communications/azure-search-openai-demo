"""Extract embedded images from PPTX and DOCX files and merge them into parsed Pages.

Uses python-pptx and python-docx to iterate through shapes/paragraphs,
extract raster images, and create ImageOnPage objects that the downstream
figure-processing pipeline (GPT-4o description, blob upload, etc.) already handles.
"""

import hashlib
import io
import logging
import os

from PIL import Image

from .page import ImageOnPage, Page

logger = logging.getLogger(__name__)

# Filtering thresholds
_MIN_IMAGE_BYTES = 2048  # 2 KB — skip tiny icons/spacers
_MIN_IMAGE_DIMENSION = 50  # px — skip bullets and decorations
_SKIP_FORMATS = {"emf", "wmf", "x-emf", "x-wmf"}
_CONTEXT_TEXT_MAX_CHARS = 2000

# Map common Office image content-types to extensions
_MIME_TO_EXT = {
    "image/png": ".png",
    "image/jpeg": ".jpg",
    "image/gif": ".gif",
    "image/tiff": ".tiff",
    "image/bmp": ".bmp",
}


def _should_skip_image(image_bytes: bytes, content_type: str) -> bool:
    """Return True if the image should be filtered out."""
    subtype = content_type.split("/")[-1].lower() if "/" in content_type else content_type.lower()
    if subtype in _SKIP_FORMATS:
        return True

    if len(image_bytes) < _MIN_IMAGE_BYTES:
        return True

    try:
        with Image.open(io.BytesIO(image_bytes)) as img:
            w, h = img.size
            if w < _MIN_IMAGE_DIMENSION or h < _MIN_IMAGE_DIMENSION:
                return True
    except Exception:
        pass

    return False


def _find_page_for_offset(char_offset: int, pages: list[Page]) -> int:
    """Find which page a character offset belongs to."""
    for i in range(len(pages) - 1, -1, -1):
        if char_offset >= pages[i].offset:
            return pages[i].page_num
    return pages[0].page_num if pages else 0


def _extract_pptx_images(document_bytes: bytes, filename: str) -> list[ImageOnPage]:
    """Extract images from a PPTX file, one per slide shape."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(io.BytesIO(document_bytes))
    images: list[ImageOnPage] = []
    seen_hashes: set[str] = set()
    img_counter = 0

    for slide_idx, slide in enumerate(prs.slides):
        # Collect slide-level context: title and full text
        context_title: str | None = None
        if slide.shapes.title is not None:
            context_title = slide.shapes.title.text.strip() or None

        text_parts: list[str] = []
        for s in slide.shapes:
            if s.has_text_frame:
                for paragraph in s.text_frame.paragraphs:
                    t = paragraph.text.strip()
                    if t:
                        text_parts.append(t)
        context_text = "\n".join(text_parts)[:_CONTEXT_TEXT_MAX_CHARS] or None

        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue

            image_blob = shape.image.blob
            content_type = shape.image.content_type or "image/png"

            if _should_skip_image(image_blob, content_type):
                continue

            img_hash = hashlib.sha256(image_blob).hexdigest()
            if img_hash in seen_hashes:
                continue
            seen_hashes.add(img_hash)

            # Extract alt text from the shape's cNvPr element (descr attribute)
            alt_text: str | None = None
            try:
                nsmap = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
                cNvPr = shape._element.find(".//p:cNvPr", nsmap)
                if cNvPr is not None:
                    descr = cNvPr.get("descr")
                    if descr and descr.strip():
                        alt_text = descr.strip()
            except Exception:
                pass

            img_counter += 1
            figure_id = f"img_{img_counter}"
            ext = _MIME_TO_EXT.get(content_type, ".png")
            base = os.path.splitext(filename)[0]
            img_filename = f"{base}_{figure_id}{ext}"
            placeholder = f'<figure id="{figure_id}"></figure>'

            images.append(
                ImageOnPage(
                    bytes=image_blob,
                    bbox=(0, 0, 0, 0),
                    filename=img_filename,
                    figure_id=figure_id,
                    page_num=slide_idx,
                    placeholder=placeholder,
                    mime_type=content_type,
                    context_title=context_title,
                    context_text=context_text,
                    alt_text=alt_text,
                )
            )

    return images


def _extract_docx_images(document_bytes: bytes, filename: str, pages: list[Page]) -> list[ImageOnPage]:
    """Extract inline images from a DOCX file, resolving page numbers from parsed pages."""
    from docx import Document
    from docx.opc.constants import RELATIONSHIP_TYPE as RT

    _HEADING_LEVELS = {
        "Heading 1": "# ",
        "Heading 2": "## ",
        "Heading 3": "### ",
        "Heading 4": "#### ",
    }

    doc = Document(io.BytesIO(document_bytes))
    images: list[ImageOnPage] = []
    seen_hashes: set[str] = set()
    img_counter = 0

    # Build cumulative character offsets per paragraph to match against Page.offset ranges
    paragraph_offsets: list[int] = []
    cumulative = 0
    for para in doc.paragraphs:
        paragraph_offsets.append(cumulative)
        cumulative += len(para.text) + 1  # +1 for implicit newline

    # Build a map from page_num -> page text for context_text lookup
    page_text_map: dict[int, str] = {p.page_num: p.text[:_CONTEXT_TEXT_MAX_CHARS] for p in pages}

    # Track the nearest heading above each paragraph
    current_heading: str | None = None

    for para_idx, para in enumerate(doc.paragraphs):
        # Update heading tracker
        style_name = para.style.name if para.style else ""
        if style_name in _HEADING_LEVELS:
            heading_text = para.text.strip()
            if heading_text:
                current_heading = heading_text

        blips = para._element.findall(
            ".//{http://schemas.openxmlformats.org/drawingml/2006/main}blip"
        )
        if not blips:
            continue

        # Collect all docPr elements in this paragraph to extract alt text per drawing
        _WP_NS = "http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing"
        _A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
        drawings = para._element.findall(
            ".//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}drawing"
        )
        # Build a map: embed_id -> alt_text by walking each drawing's docPr + blip
        blip_alt_map: dict[str, str] = {}
        for drawing in drawings:
            doc_prs = drawing.findall(f".//{{{_WP_NS}}}docPr")
            drawing_blips = drawing.findall(f".//{{{_A_NS}}}blip")
            if doc_prs and drawing_blips:
                descr = doc_prs[0].get("descr", "").strip()
                for db in drawing_blips:
                    embed = db.get(
                        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
                    )
                    if embed and descr:
                        blip_alt_map[embed] = descr

        for blip in blips:
            embed_id = blip.get(
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
            )
            if not embed_id:
                continue

            try:
                rel = para.part.rels[embed_id]
            except KeyError:
                continue

            if rel.reltype != RT.IMAGE:
                continue

            image_blob = rel.target_part.blob
            content_type = rel.target_part.content_type or "image/png"

            if _should_skip_image(image_blob, content_type):
                continue

            img_hash = hashlib.sha256(image_blob).hexdigest()
            if img_hash in seen_hashes:
                continue
            seen_hashes.add(img_hash)

            img_counter += 1
            figure_id = f"img_{img_counter}"
            ext = _MIME_TO_EXT.get(content_type, ".png")
            base = os.path.splitext(filename)[0]
            img_filename = f"{base}_{figure_id}{ext}"
            placeholder = f'<figure id="{figure_id}"></figure>'

            para_offset = paragraph_offsets[para_idx]
            page_num = _find_page_for_offset(para_offset, pages) if pages else 0

            # Context fields
            context_title = current_heading
            context_text = page_text_map.get(page_num)
            alt_text = blip_alt_map.get(embed_id) or None

            images.append(
                ImageOnPage(
                    bytes=image_blob,
                    bbox=(0, 0, 0, 0),
                    filename=img_filename,
                    figure_id=figure_id,
                    page_num=page_num,
                    placeholder=placeholder,
                    mime_type=content_type,
                    context_title=context_title,
                    context_text=context_text,
                    alt_text=alt_text,
                )
            )

    return images


def extract_and_merge_office_images(filename: str, document_bytes: bytes, pages: list[Page]) -> list[Page]:
    """Extract embedded images from PPTX/DOCX and merge them into existing pages.

    Creates ImageOnPage objects with placeholders, appends placeholders to page text,
    and adds images to the corresponding Page.images list. The downstream pipeline
    (GPT-4o description, blob upload, placeholder replacement) handles the rest.

    Args:
        filename: Original document filename (e.g. "slides.pptx")
        document_bytes: Raw file bytes
        pages: Parsed pages from Document Intelligence

    Returns:
        The same pages list, mutated with extracted images.
    """
    ext = os.path.splitext(filename)[1].lower()

    if ext == ".pptx":
        images = _extract_pptx_images(document_bytes, filename)
    elif ext == ".docx":
        images = _extract_docx_images(document_bytes, filename, pages)
    else:
        return pages

    if not images:
        logger.info("No extractable images found in %s", filename)
        return pages

    logger.info("Extracted %d images from %s", len(images), filename)

    # Build a lookup of page_num → Page for fast assignment
    page_map: dict[int, Page] = {p.page_num: p for p in pages}

    for image in images:
        target_page = page_map.get(image.page_num)
        if target_page is None:
            # Fall back to the last page if the image's page_num doesn't match
            target_page = pages[-1] if pages else None
            if target_page is None:
                continue

        # Append placeholder to end of page text
        target_page.text = target_page.text.rstrip() + "\n" + image.placeholder
        target_page.images.append(image)

    return pages
