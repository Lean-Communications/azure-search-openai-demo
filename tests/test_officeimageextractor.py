"""Tests for Office image extraction with context fields (PPTX & DOCX)."""

import io

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from prepdocslib.officeimageextractor import (
    _CONTEXT_TEXT_MAX_CHARS,
    _extract_docx_images,
    _extract_pptx_images,
)
from prepdocslib.page import Page


def _make_large_test_png() -> bytes:
    """Create a noisy PNG that is guaranteed to pass the 2 KB byte-size filter.

    Solid-color PNGs compress to ~300-600 bytes, which is below the 2048-byte
    minimum threshold.  A random-noise image compresses much larger.
    """
    import random

    rng = random.Random(42)  # deterministic for reproducibility
    pixels = bytes([rng.randint(0, 255) for _ in range(200 * 200 * 3)])
    img = Image.frombytes("RGB", (200, 200), pixels)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    data = buf.getvalue()
    assert len(data) >= 2048, f"Test image only {len(data)} bytes — too small for filters"
    return data


def _build_pptx_bytes(
    title_text: str | None = "Test Title",
    body_text: str = "Some body text",
    include_picture: bool = True,
) -> bytes:
    """Build a minimal PPTX in memory with one slide."""
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # Layout with title and content
    slide = prs.slides.add_slide(slide_layout)

    if title_text is not None and slide.shapes.title is not None:
        slide.shapes.title.text = title_text

    # Add a text box
    from pptx.util import Emu

    left = Inches(1)
    top = Inches(2)
    width = Inches(4)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.text = body_text

    if include_picture:
        img_bytes = _make_large_test_png()
        img_stream = io.BytesIO(img_bytes)
        slide.shapes.add_picture(img_stream, Inches(1), Inches(3), Inches(2), Inches(2))

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _build_pptx_no_title_bytes() -> bytes:
    """Build a PPTX with a blank slide layout (no title placeholder)."""
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Blank layout - no title
    slide = prs.slides.add_slide(slide_layout)

    # Add only a picture, no title shape
    img_bytes = _make_large_test_png()
    img_stream = io.BytesIO(img_bytes)
    slide.shapes.add_picture(img_stream, Inches(1), Inches(1), Inches(2), Inches(2))

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _build_docx_bytes(
    heading_text: str = "Section Header",
    body_text: str = "Some paragraph text under the heading.",
    include_image: bool = True,
) -> bytes:
    """Build a minimal DOCX in memory with a heading, body, and optional inline image."""
    from docx import Document
    from docx.shared import Inches as DocxInches

    doc = Document()
    doc.add_heading(heading_text, level=1)
    doc.add_paragraph(body_text)

    if include_image:
        img_bytes = _make_large_test_png()
        img_stream = io.BytesIO(img_bytes)
        doc.add_picture(img_stream, DocxInches(2))

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# PPTX tests
# ---------------------------------------------------------------------------


class TestPptxContextFields:
    def test_pptx_context_fields(self):
        """Images extracted from PPTX have slide title and context text."""
        pptx_bytes = _build_pptx_bytes(title_text="Test Title", body_text="Some body text")
        images = _extract_pptx_images(pptx_bytes, "slides.pptx")

        assert len(images) >= 1
        img = images[0]
        assert img.context_title == "Test Title"
        assert img.context_text is not None
        assert "Some body text" in img.context_text

    def test_pptx_no_title_slide(self):
        """Slide with no title shape produces context_title=None."""
        pptx_bytes = _build_pptx_no_title_bytes()
        images = _extract_pptx_images(pptx_bytes, "slides.pptx")

        assert len(images) >= 1
        img = images[0]
        assert img.context_title is None

    def test_pptx_context_text_includes_title(self):
        """The context_text includes the title text (since it is a text shape on the slide)."""
        pptx_bytes = _build_pptx_bytes(title_text="Slide Title", body_text="Body content here")
        images = _extract_pptx_images(pptx_bytes, "slides.pptx")

        assert len(images) >= 1
        # The title shape is also a text shape, so context_text should include it
        assert images[0].context_text is not None
        assert "Slide Title" in images[0].context_text
        assert "Body content here" in images[0].context_text


# ---------------------------------------------------------------------------
# DOCX tests
# ---------------------------------------------------------------------------


class TestDocxContextFields:
    def test_docx_context_fields(self):
        """Images extracted from DOCX have the nearest heading as context_title."""
        docx_bytes = _build_docx_bytes(heading_text="Section Header", body_text="Body text.")

        # Create a Page that covers the whole document text
        pages = [Page(page_num=0, offset=0, text="Section Header\nBody text.\n")]

        images = _extract_docx_images(docx_bytes, "document.docx", pages)

        assert len(images) >= 1
        img = images[0]
        assert img.context_title == "Section Header"

    def test_docx_context_text_from_pages(self):
        """context_text comes from the pages list, not raw paragraph text."""
        docx_bytes = _build_docx_bytes(heading_text="Heading", body_text="Paragraph text.")

        page_text = "Heading\nParagraph text.\nMore text from DI parsing."
        pages = [Page(page_num=0, offset=0, text=page_text)]

        images = _extract_docx_images(docx_bytes, "document.docx", pages)

        assert len(images) >= 1
        assert images[0].context_text == page_text

    def test_docx_no_heading(self):
        """DOCX with no heading paragraphs before image gives context_title=None."""
        from docx import Document
        from docx.shared import Inches as DocxInches

        doc = Document()
        doc.add_paragraph("Just a plain paragraph, no heading.")
        img_bytes = _make_large_test_png()
        doc.add_picture(io.BytesIO(img_bytes), DocxInches(2))

        buf = io.BytesIO()
        doc.save(buf)
        docx_bytes = buf.getvalue()

        pages = [Page(page_num=0, offset=0, text="Just a plain paragraph, no heading.\n")]
        images = _extract_docx_images(docx_bytes, "document.docx", pages)

        assert len(images) >= 1
        assert images[0].context_title is None


# ---------------------------------------------------------------------------
# Truncation test
# ---------------------------------------------------------------------------


class TestContextTextTruncation:
    def test_context_text_truncation_pptx(self):
        """context_text is truncated to _CONTEXT_TEXT_MAX_CHARS for PPTX."""
        long_text = "A" * (_CONTEXT_TEXT_MAX_CHARS + 500)
        pptx_bytes = _build_pptx_bytes(title_text="Title", body_text=long_text)
        images = _extract_pptx_images(pptx_bytes, "slides.pptx")

        assert len(images) >= 1
        assert images[0].context_text is not None
        assert len(images[0].context_text) <= _CONTEXT_TEXT_MAX_CHARS

    def test_context_text_truncation_docx(self):
        """context_text is truncated to _CONTEXT_TEXT_MAX_CHARS for DOCX."""
        long_page_text = "B" * (_CONTEXT_TEXT_MAX_CHARS + 500)
        docx_bytes = _build_docx_bytes(heading_text="H", body_text="text")

        pages = [Page(page_num=0, offset=0, text=long_page_text)]
        images = _extract_docx_images(docx_bytes, "document.docx", pages)

        assert len(images) >= 1
        assert images[0].context_text is not None
        assert len(images[0].context_text) <= _CONTEXT_TEXT_MAX_CHARS
