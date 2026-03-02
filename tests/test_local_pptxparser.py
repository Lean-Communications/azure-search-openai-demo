"""Tests for LocalPptxParser text extraction."""

import asyncio
import io
import random

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from prepdocslib.page import Page
from prepdocslib.pdfparser import LocalPptxParser


def _make_large_test_png() -> bytes:
    """Create a noisy PNG that is guaranteed to pass the 2 KB byte-size filter.

    Solid-color PNGs compress to ~300-600 bytes, which is below the 2048-byte
    minimum threshold.  A random-noise image compresses much larger.
    """
    rng = random.Random(42)  # deterministic for reproducibility
    pixels = bytes([rng.randint(0, 255) for _ in range(200 * 200 * 3)])
    img = Image.frombytes("RGB", (200, 200), pixels)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    data = buf.getvalue()
    assert len(data) >= 2048, f"Test image only {len(data)} bytes — too small for filters"
    return data


def _build_pptx(slides: list[dict]) -> bytes:
    """Build a PPTX in memory from a list of slide dicts.

    Each dict may contain:
        title (str | None): Title text for the slide.
        body (str | None): Body text added via a textbox.
        notes (str | None): Speaker notes text.
        include_picture (bool): Whether to embed a test PNG image.
        table (list[list[str]] | None): 2-D list of cell values for a table.
    """
    prs = Presentation()

    for slide_spec in slides:
        title = slide_spec.get("title")
        body = slide_spec.get("body")
        notes = slide_spec.get("notes")
        include_picture = slide_spec.get("include_picture", False)
        table_data = slide_spec.get("table")

        # Choose layout: blank (6) when nothing is provided, title+content (1) otherwise
        is_empty = title is None and body is None and notes is None and not include_picture and table_data is None
        layout_index = 6 if is_empty else 1
        slide_layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(slide_layout)

        # Title
        if title is not None and slide.shapes.title is not None:
            slide.shapes.title.text = title

        # Body text via textbox
        if body is not None:
            left = Inches(1)
            top = Inches(2)
            width = Inches(4)
            height = Inches(1)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            txBox.text_frame.text = body

        # Speaker notes
        if notes is not None:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

        # Picture
        if include_picture:
            img_bytes = _make_large_test_png()
            img_stream = io.BytesIO(img_bytes)
            slide.shapes.add_picture(img_stream, Inches(1), Inches(3), Inches(2), Inches(2))

        # Table
        if table_data is not None:
            rows = len(table_data)
            cols = len(table_data[0]) if rows > 0 else 0
            tbl = slide.shapes.add_table(rows, cols, Inches(1), Inches(4), Inches(6), Inches(2)).table
            for r_idx, row in enumerate(table_data):
                for c_idx, cell_text in enumerate(row):
                    tbl.cell(r_idx, c_idx).text = cell_text

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _parse_pptx_sync(pptx_bytes: bytes, filename: str = "test.pptx") -> list[Page]:
    """Run the async LocalPptxParser.parse synchronously and return a list of Pages."""
    stream = io.BytesIO(pptx_bytes)
    stream.name = filename
    parser = LocalPptxParser()

    async def _collect() -> list[Page]:
        pages: list[Page] = []
        async for page in parser.parse(stream):
            pages.append(page)
        return pages

    return asyncio.run(_collect())


class TestLocalPptxParserText:
    """Tests for basic slide text extraction via LocalPptxParser."""

    def test_single_slide_title_and_body(self):
        """Single slide yields one Page with '# My Title' and body text."""
        pptx_bytes = _build_pptx([{"title": "My Title", "body": "Hello world"}])
        pages = _parse_pptx_sync(pptx_bytes)

        assert len(pages) == 1
        page = pages[0]
        assert page.page_num == 0
        assert "# My Title" in page.text
        assert "Hello world" in page.text

    def test_multiple_slides(self):
        """3 slides yield 3 Pages with correct page_num (0, 1, 2) and titles."""
        pptx_bytes = _build_pptx([
            {"title": "First", "body": "a"},
            {"title": "Second", "body": "b"},
            {"title": "Third", "body": "c"},
        ])
        pages = _parse_pptx_sync(pptx_bytes)

        assert len(pages) == 3
        for i, expected_title in enumerate(["First", "Second", "Third"]):
            assert pages[i].page_num == i
            assert f"# {expected_title}" in pages[i].text

    def test_offsets_are_cumulative(self):
        """page[1].offset == len(page[0].text)."""
        pptx_bytes = _build_pptx([
            {"title": "Slide One", "body": "Content A"},
            {"title": "Slide Two", "body": "Content B"},
        ])
        pages = _parse_pptx_sync(pptx_bytes)

        assert len(pages) == 2
        assert pages[0].offset == 0
        assert pages[1].offset == len(pages[0].text)

    def test_speaker_notes(self):
        """Notes text appears after 'Notes:' separator."""
        pptx_bytes = _build_pptx([{"title": "Titled", "notes": "Remember this"}])
        pages = _parse_pptx_sync(pptx_bytes)

        assert len(pages) == 1
        assert "Notes:" in pages[0].text
        assert "Remember this" in pages[0].text

    def test_table_extraction(self):
        """Tables rendered as pipe-delimited rows ('Name | Age')."""
        pptx_bytes = _build_pptx([{
            "title": "Data",
            "table": [
                ["Name", "Age"],
                ["Alice", "30"],
            ],
        }])
        pages = _parse_pptx_sync(pptx_bytes)

        assert len(pages) == 1
        assert "Name | Age" in pages[0].text
        assert "Alice | 30" in pages[0].text

    def test_empty_slide(self):
        """Blank slide layout yields a Page with page_num=0."""
        pptx_bytes = _build_pptx([{}])
        pages = _parse_pptx_sync(pptx_bytes)

        assert len(pages) == 1
        assert pages[0].page_num == 0


class TestLocalPptxParserImages:
    def test_images_attached_to_pages(self):
        """Slides with pictures have ImageOnPage objects attached."""
        pptx_bytes = _build_pptx([{"title": "Img Slide", "body": "Text", "include_picture": True}])
        pages = _parse_pptx_sync(pptx_bytes)

        assert len(pages) == 1
        assert len(pages[0].images) >= 1
        img = pages[0].images[0]
        assert img.page_num == 0
        assert img.context_title == "Img Slide"
        assert len(img.bytes) > 0

    def test_image_placeholder_in_text(self):
        """Image placeholders are appended to the page text."""
        pptx_bytes = _build_pptx([{"title": "Pic", "body": "Words", "include_picture": True}])
        pages = _parse_pptx_sync(pptx_bytes)

        assert len(pages) == 1
        assert "<figure" in pages[0].text

    def test_slide_without_image_has_no_images(self):
        """Slides without pictures have empty images list."""
        pptx_bytes = _build_pptx([
            {"title": "No Pic", "body": "Just text"},
            {"title": "Has Pic", "body": "With image", "include_picture": True},
        ])
        pages = _parse_pptx_sync(pptx_bytes)

        assert len(pages) == 2
        assert len(pages[0].images) == 0
        assert len(pages[1].images) >= 1
